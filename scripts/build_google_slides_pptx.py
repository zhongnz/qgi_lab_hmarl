#!/usr/bin/env python3
"""Build a cleaner PowerPoint deck that imports well into Google Slides."""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "docs" / "reports"
RUNS = ROOT / "runs"
OUT = REPORTS / "2026-03-31_hmarl_project_update_google_slides_import.pptx"

V3_RUN = RUNS / "local_full_train_2026-03-11_transit_rebalanced_v3"
EVAL_TRACE_CSV = V3_RUN / "eval_trace.csv"
ROLE_DIAG_IMG = REPORTS / "2026-03-31_role_reward_diagnostics.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(244, 239, 230)
PAPER = RGBColor(255, 250, 242)
INK = RGBColor(20, 35, 54)
MUTED = RGBColor(84, 98, 117)
NAVY = RGBColor(20, 56, 92)
NAVY_2 = RGBColor(34, 73, 110)
TEAL = RGBColor(43, 122, 120)
TEAL_SOFT = RGBColor(235, 247, 245)
ORANGE = RGBColor(214, 126, 63)
ORANGE_SOFT = RGBColor(253, 243, 232)
SAND = RGBColor(231, 216, 189)
LINE = RGBColor(221, 209, 188)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(57, 130, 93)
RED = RGBColor(183, 79, 66)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Courier New"


def _set_slide_bg(slide, color: RGBColor = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_shape(shape, fill_rgb: RGBColor, line_rgb: RGBColor | None = LINE) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb


def _add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    font_size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT_BODY,
    word_wrap: bool = True,
) :
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _add_title_band(slide, title: str, eyebrow: str, page: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.88))
    _style_shape(band, NAVY, None)
    _add_textbox(slide, 0.45, 0.18, 2.1, 0.25, eyebrow.upper(), font_size=12, color=SAND, bold=True)
    _add_textbox(slide, 0.45, 0.38, 10.7, 0.34, title, font_size=24, color=PAPER, bold=True, font_name=FONT_HEAD)
    _add_textbox(slide, 12.35, 0.2, 0.45, 0.22, page, font_size=16, color=SAND, bold=True, align=PP_ALIGN.RIGHT)


def _add_card(slide, x: float, y: float, w: float, h: float, *, fill: RGBColor = WHITE, title: str | None = None, title_color: RGBColor = NAVY):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _style_shape(card, fill, LINE)
    if title:
        _add_textbox(slide, x + 0.22, y + 0.18, w - 0.4, 0.28, title, font_size=16, color=title_color, bold=True, font_name=FONT_HEAD)
    return card


def _add_metric_card(slide, x: float, y: float, w: float, h: float, label: str, value: str, note: str, *, fill: RGBColor = WHITE, value_color: RGBColor = NAVY) -> None:
    _add_card(slide, x, y, w, h, fill=fill)
    _add_textbox(slide, x + 0.22, y + 0.18, w - 0.4, 0.18, label.upper(), font_size=10, color=MUTED, bold=True)
    _add_textbox(slide, x + 0.22, y + 0.48, w - 0.4, 0.36, value, font_size=22, color=value_color, bold=True, font_name=FONT_HEAD)
    _add_textbox(slide, x + 0.22, y + 0.93, w - 0.4, 0.3, note, font_size=11, color=MUTED)


def _add_bullets(box, items: list[str], *, font_size: int = 18, color: RGBColor = INK, font_name: str = FONT_BODY, first_indent: int | None = None) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        p.space_after = Pt(6)
        if first_indent is not None:
            p.level = 0
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def _read_v3_metrics() -> dict[str, float]:
    data = json.loads((V3_RUN / "eval_result.json").read_text())["mean"]
    return {
        "total_reward": float(data["total_reward"]),
        "on_time_rate": float(data["on_time_rate"]),
        "completed_arrivals": float(data["completed_arrivals"]),
        "total_vessels_served": float(data["total_vessels_served"]),
        "dock_utilization": float(data["dock_utilization"]),
        "total_ops_cost_usd": float(data["total_ops_cost_usd"]),
    }


def _ensure_role_reward_plot() -> Path:
    if ROLE_DIAG_IMG.exists() and ROLE_DIAG_IMG.stat().st_mtime >= EVAL_TRACE_CSV.stat().st_mtime:
        return ROLE_DIAG_IMG

    with EVAL_TRACE_CSV.open() as fh:
        rows = list(csv.DictReader(fh))

    t = [float(row["t"]) for row in rows]
    series = [
        ("Vessel reward", [float(row["avg_vessel_reward"]) for row in rows], (43 / 255, 122 / 255, 120 / 255)),
        ("Port reward", [float(row["avg_port_reward"]) for row in rows], (214 / 255, 126 / 255, 63 / 255)),
        ("Coordinator reward", [float(row["coordinator_reward"]) for row in rows], (105 / 255, 110 / 255, 196 / 255)),
    ]

    fig, axes = plt.subplots(1, 3, figsize=(10.8, 3.0), dpi=180)
    fig.patch.set_facecolor((250 / 255, 248 / 255, 244 / 255))
    for ax, (title, values, color) in zip(axes, series):
        ax.plot(t, values, color=color, linewidth=2.2)
        ax.set_title(title, fontsize=12)
        ax.set_xlabel("t", fontsize=10)
        ax.grid(True, alpha=0.25)
        ax.tick_params(labelsize=9)
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
    axes[0].set_ylabel("reward", fontsize=10)
    fig.suptitle("Per-step reward diagnostics by role", fontsize=16, y=1.02)
    fig.tight_layout()
    fig.savefig(ROLE_DIAG_IMG, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return ROLE_DIAG_IMG


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(0.8), Inches(3.25), Inches(2.15))
    _style_shape(accent, RGBColor(229, 238, 245), None)

    accent2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(3.35), Inches(2.7), Inches(1.45))
    _style_shape(accent2, ORANGE_SOFT, None)

    _add_textbox(
        slide,
        0.7,
        0.95,
        8.0,
        1.55,
        "Hierarchical Multi-Agent Reinforcement Learning for Congestion-Aware Maritime Scheduling",
        font_size=26,
        color=INK,
        bold=True,
        font_name=FONT_HEAD,
    )
    _add_textbox(
        slide,
        0.75,
        2.55,
        7.5,
        0.75,
        "Project update: simulator background, reward redesign, continuous + ground_truth final run path, current progress, and next steps.",
        font_size=16,
        color=MUTED,
    )

    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.5), Inches(7.65), Inches(1.25))
    _style_shape(quote, NAVY, None)
    _add_textbox(
        slide,
        1.0,
        3.8,
        7.15,
        0.8,
        "We now have a clean final path for the project: continuous scheduling, reliable forecasts, strong diagnostics, and a reproducible local run plan.",
        font_size=17,
        color=PAPER,
    )

    _add_textbox(slide, 9.3, 1.08, 2.55, 0.16, "SPRING 2026", font_size=10, color=MUTED, bold=True)
    _add_textbox(slide, 9.3, 1.42, 2.55, 0.26, "Independent study", font_size=18, color=NAVY, bold=True, font_name=FONT_HEAD)
    _add_textbox(slide, 9.3, 1.84, 2.55, 0.42, "Supervisor:\nProf. Amine Aboussalah", font_size=13, color=INK)

    _add_textbox(slide, 9.7, 3.68, 2.0, 0.14, "MARCH 31 UPDATE", font_size=9, color=MUTED, bold=True)
    _add_textbox(slide, 9.7, 4.0, 2.0, 0.25, "Final path:\ncontinuous + GT", font_size=13, color=NAVY, bold=True)

    _add_metric_card(slide, 0.75, 5.55, 3.1, 1.35, "Research focus", "HMARL", "Coordinator, vessel, and port agents")
    _add_metric_card(slide, 4.05, 5.55, 3.35, 1.35, "Current setup", "Cont. + GT", "Continuous environment with reliable forecasts")
    _add_metric_card(slide, 7.65, 5.55, 2.95, 1.35, "Format", "12 slides", "Presentation-ready update deck")


def _context_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Why this project matters", "Context", "02")

    _add_card(slide, 0.65, 1.4, 6.2, 5.65, fill=RGBColor(238, 246, 251), title="Operational context")
    _add_bullets(
        _add_textbox(slide, 0.95, 1.95, 5.55, 4.6, "", font_size=1),
        [
            "Ports have limited berth capacity and nonzero service time.",
            "Vessels trade off speed, fuel burn, emissions, and lateness.",
            "Poor local choices can create downstream congestion and idle capacity elsewhere.",
            "The system is naturally multi-agent, asynchronous, and resource-constrained.",
        ],
        font_size=19,
    )

    _add_card(slide, 7.1, 1.4, 2.65, 2.55, fill=WHITE, title="Congestion")
    _add_textbox(slide, 7.35, 1.95, 2.15, 1.45, "Too many arrivals to one port can create queueing and delay at the whole network level.", font_size=14, color=INK)

    _add_card(slide, 9.95, 1.4, 2.65, 2.55, fill=WHITE, title="Cost")
    _add_textbox(slide, 10.2, 1.95, 2.15, 1.45, "Faster sailing can improve timing but increases fuel use, emissions, and operating cost.", font_size=14, color=INK)

    _add_card(slide, 7.1, 4.2, 2.65, 2.55, fill=WHITE, title="Coordination")
    _add_textbox(slide, 7.35, 4.75, 2.15, 1.45, "Ports and vessels need global guidance, not only local reactive behavior.", font_size=14, color=INK)

    _add_card(slide, 9.95, 4.2, 2.65, 2.55, fill=WHITE, title="Research value")
    _add_textbox(slide, 10.2, 4.75, 2.15, 1.45, "This is a strong testbed for hierarchical coordination under delayed information.", font_size=14, color=INK)


def _rq_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "What we are trying to learn", "Questions", "03")

    hero = _add_card(slide, 0.65, 1.4, 12.0, 1.0, fill=NAVY, title="Umbrella question", title_color=SAND)
    hero.line.fill.background()
    _add_textbox(
        slide,
        0.95,
        1.88,
        11.4,
        0.28,
        "Can hierarchical MARL improve maritime scheduling under congestion and resource constraints?",
        font_size=18,
        color=PAPER,
        bold=True,
    )

    questions = [
        ("RQ1", "How can heterogeneous agents coordinate using shared congestion forecasts?"),
        ("RQ2", "Does proactive coordination improve over independent and reactive baselines?"),
        ("RQ3", "Which forecast horizons and sharing strategies maximize decision quality?"),
        ("RQ4", "How do coordination improvements affect fuel, delay, and carbon cost?"),
    ]
    positions = [(0.65, 2.75), (6.65, 2.75), (0.65, 4.8), (6.65, 4.8)]
    for (rq, text), (x, y) in zip(questions, positions):
        _add_card(slide, x, y, 5.6, 1.6, fill=WHITE, title=rq)
        _add_textbox(slide, x + 0.25, y + 0.58, 5.0, 0.75, text, font_size=15, color=INK)

    _add_textbox(
        slide,
        0.85,
        6.7,
        11.7,
        0.38,
        "The project is organized around a chain of related questions: coordination, baselines, forecast design, and economics.",
        font_size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def _environment_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Current simulator setup", "Environment", "04")

    # Left stack of configuration cards
    cards = [
        ("Topology", "5 ports, 8 vessels, 3 docks per port"),
        ("Time scale", "1 hour per step and 6 hour default service time"),
        ("Travel", "Synthetic port-distance matrix with routes from 22 to 98 nautical miles"),
        ("Fuel", "Default fuel is 100, with refueling only after real service completion"),
    ]
    y = 1.45
    for title, text in cards:
        _add_card(slide, 0.65, y, 4.0, 1.2, fill=WHITE, title=title)
        _add_textbox(slide, 0.9, y + 0.52, 3.45, 0.45, text, font_size=13, color=INK)
        y += 1.38

    _add_card(slide, 4.95, 1.45, 7.7, 5.85, fill=RGBColor(250, 248, 244), title="Synthetic port distance matrix (nautical miles)")
    _add_textbox(slide, 5.2, 1.9, 7.15, 0.5, "Rows and columns are ports P0–P4; the diagonal is 0, and the matrix is symmetric.", font_size=13, color=MUTED)
    _add_textbox(slide, 5.2, 2.25, 7.15, 0.6, "The values are synthetic and scaled so multiple voyages can complete within one rollout.", font_size=13, color=MUTED)

    labels = ["", "P0", "P1", "P2", "P3", "P4"]
    matrix = [
        ["P0", "0", "84", "22", "34", "78"],
        ["P1", "84", "0", "98", "61", "54"],
        ["P2", "22", "98", "0", "55", "59"],
        ["P3", "34", "61", "55", "0", "82"],
        ["P4", "78", "54", "59", "82", "0"],
    ]
    cell_w = 0.95
    cell_h = 0.5
    start_x = 5.35
    start_y = 3.05
    for col, label in enumerate(labels):
        _add_card(slide, start_x + col * cell_w, start_y, cell_w, cell_h, fill=RGBColor(228, 236, 244), title=None)
        _add_textbox(slide, start_x + col * cell_w + 0.03, start_y + 0.14, cell_w - 0.06, 0.16, label, font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    for row_idx, row in enumerate(matrix, start=1):
        for col, value in enumerate(row):
            fill = WHITE if col else RGBColor(247, 244, 239)
            _add_card(slide, start_x + col * cell_w, start_y + row_idx * cell_h, cell_w, cell_h, fill=fill, title=None)
            _add_textbox(
                slide,
                start_x + col * cell_w + 0.03,
                start_y + row_idx * cell_h + 0.14,
                cell_w - 0.06,
                0.16,
                value,
                font_size=12,
                color=INK if col else NAVY,
                bold=(col == 0),
                align=PP_ALIGN.CENTER,
            )


def _experiment_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "How the current baseline is trained and evaluated", "Experiment", "05")

    _add_card(slide, 0.65, 1.45, 5.85, 5.8, fill=RGBColor(245, 250, 252), title="Main operating baseline")
    left_box = _add_textbox(slide, 0.95, 1.95, 5.2, 4.95, "", font_size=1)
    _add_bullets(
        left_box,
        [
            "Algorithm: MAPPO with centralized training and decentralized execution",
            "Parameter sharing across vessels and across ports",
            "100 training iterations and rollout_length = 64",
            "Environment horizon = 69",
            "Seed = 42",
            "Local CPU run for the reported baseline",
        ],
        font_size=17,
    )

    _add_card(slide, 6.8, 1.45, 5.85, 2.75, fill=ORANGE_SOFT, title="Final completion path")
    right_top = _add_textbox(slide, 7.1, 1.95, 5.2, 1.8, "", font_size=1)
    _add_bullets(
        right_top,
        [
            "continuous episodes",
            "ground_truth forecasts",
            "one artifact run plus one five-seed run",
            "designed to isolate control quality while finishing the project",
        ],
        font_size=15,
    )

    _add_card(slide, 6.8, 4.5, 5.85, 2.75, fill=TEAL_SOFT, title="Primary evaluation metrics")
    right_bottom = _add_textbox(slide, 7.1, 5.0, 5.2, 1.8, "", font_size=1)
    _add_bullets(
        right_bottom,
        [
            "On-time rate",
            "Completed arrivals",
            "Port service events",
            "Dock utilization",
            "Total operating cost",
        ],
        font_size=16,
    )


def _architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Three decision layers with delayed coordination", "Architecture", "06")

    items = [
        ("Fleet coordinator", "Strategic destination guidance\nSlowest decision cadence", RGBColor(245, 248, 252)),
        ("Vessel agents", "Speed control\nArrival-slot requests", WHITE),
        ("Port agents", "Request acceptance\nService and berth decisions", RGBColor(252, 248, 242)),
    ]
    x_positions = [0.9, 4.65, 8.4]
    for (title, text, fill), x in zip(items, x_positions):
        _add_card(slide, x, 2.25, 3.0, 2.2, fill=fill, title=title)
        _add_textbox(slide, x + 0.25, 2.95, 2.45, 1.0, text, font_size=16, color=INK)

    for x in [3.95, 7.7]:
        _add_textbox(slide, x, 3.0, 0.45, 0.5, "→", font_size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.1))
    _style_shape(callout, RGBColor(249, 237, 220), None)
    _add_textbox(
        slide,
        1.25,
        5.48,
        10.75,
        0.4,
        "The environment is intentionally asynchronous, so the policies must coordinate under delayed information and physical travel time.",
        font_size=16,
        color=INK,
        align=PP_ALIGN.CENTER,
    )


def _mechanics_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "How one simulation step works", "Mechanics", "07")

    _add_card(slide, 0.65, 1.45, 8.2, 5.9, fill=RGBColor(245, 250, 252), title="Transition kernel")
    step_titles = ["1. Deliver messages", "2. Advance vessels", "3. Process ports", "4. Compute rewards"]
    step_notes = [
        "Apply due actions and delayed communications.",
        "Update motion, fuel burn, emissions, and arrivals.",
        "Update queue, berth service, service completion, and refueling.",
        "Log diagnostics and compute role-specific rewards.",
    ]
    y = 2.0
    for title, note in zip(step_titles, step_notes):
        _add_card(slide, 0.95, y, 7.6, 0.95, fill=WHITE, title=title)
        _add_textbox(slide, 1.18, y + 0.47, 7.05, 0.28, note, font_size=13, color=INK)
        y += 1.12

    _add_card(slide, 9.15, 1.45, 3.5, 5.9, fill=RGBColor(252, 248, 242), title="Why this matters")
    _add_bullets(
        _add_textbox(slide, 9.45, 1.98, 2.95, 3.8, "", font_size=1),
        [
            "The environment is not one big update; it moves through communication, travel, berth service, and reward calculation in sequence.",
            "This is why delays, arrivals, and service completions show up at different times in the logs.",
            "The next two slides separate reward intuition from reward formulas so the audience sees both the story and the technical definition.",
        ],
        font_size=14,
    )


def _reward_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Reward structure in plain language", "Rewards", "08")

    cards = [
        (
            "Vessel reward",
            "Reward = arrival and punctuality minus travel cost",
            ["Arrival", "On-time arrival"],
            ["Fuel used", "Delay added", "CO2 emitted", "Transit time", "Schedule delay"],
            WHITE,
        ),
        (
            "Port reward",
            "Reward = responsive berth management minus congestion cost",
            ["Accepted requests", "Vessels served"],
            ["Rejected requests", "Queue waiting", "Idle docks"],
            RGBColor(248, 253, 252),
        ),
        (
            "Coordinator reward",
            "Reward = system flow and utilization minus network-wide cost",
            ["Accepted requests", "Vessels served", "Berth utilization"],
            ["Fuel", "Queue", "Idle docks", "Delay", "Schedule delay", "CO2", "Rejected requests"],
            WHITE,
        ),
    ]

    y = 1.22
    heights = [2.0, 2.0, 2.2]
    for (title, headline, rewards, penalties, fill), h in zip(cards, heights):
        _add_card(slide, 0.75, y, 12.0, h, fill=fill, title=title)
        _add_textbox(slide, 1.0, y + 0.38, 11.3, 0.2, headline, font_size=13, color=INK, bold=True)

        _add_card(slide, 1.0, y + 0.74, 5.25, h - 0.94, fill=RGBColor(242, 248, 244), title="What it rewards")
        _add_textbox(slide, 1.22, y + 1.04, 4.8, h - 1.26, ", ".join(rewards), font_size=11, color=INK)

        _add_card(slide, 6.45, y + 0.74, 5.95, h - 0.94, fill=RGBColor(252, 245, 239), title="What it penalizes")
        _add_textbox(slide, 6.67, y + 1.04, 5.45, h - 1.26, ", ".join(penalties), font_size=11, color=INK)
        y += h + 0.15

    _add_textbox(slide, 0.95, 7.18, 11.6, 0.14, "This is the audience-facing summary. The next slide shows the actual symbolic formulas used in the simulator.", font_size=10, color=MUTED, align=PP_ALIGN.CENTER)


def _reward_formula_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Actual reward formulas used in the simulator", "Rewards", "09")

    _add_card(slide, 0.75, 1.45, 12.0, 5.95, fill=RGBColor(252, 248, 242), title="Symbolic form")

    formulas = [
        (
            "Vessel reward",
            "r_V(t) = r_arr·1[arrived] + r_on·1[on_time]\n          − (w_fΔfuel + w_dΔdelay + w_eΔCO2 + w_tΔtransit + w_sΔsched)",
            WHITE,
            1.55,
            1.8,
        ),
        (
            "Port reward",
            "r_P(t) = r_acc·accepted + r_srv·served\n          − r_rej·rejected − (queue·dt + w_idle·idle_docks)",
            RGBColor(248, 253, 252),
            3.48,
            1.45,
        ),
        (
            "Coordinator reward",
            "r_C(t) = r_acc·accepted + r_srv·served + r_util·avg_occupied\n          − (w_fΔfuel + w_qavg_queue + w_iavg_idle + w_dΔdelay + w_sΔsched + w_eΔCO2 + r_rej·rejected)",
            WHITE,
            5.1,
            1.9,
        ),
    ]
    for title, formula, fill, y, h in formulas:
        _add_card(slide, 1.05, y, 11.4, h, fill=fill, title=title)
        _add_textbox(slide, 1.3, y + 0.45, 10.9, h - 0.55, formula, font_size=14, color=INK, font_name=FONT_MONO)
    _add_textbox(slide, 1.05, 7.08, 11.4, 0.2, "Δ terms are per-step increments; accepted, rejected, and served are event counts for the current step.", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)


def _changes_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "What changed since March 6", "Progress", "10")

    items = [
        ("Corrected vessel behavior", "No more collapse to constant minimum speed.\nRequested arrival times now become real schedule deadlines."),
        ("Fuel and stalling", "Fuel exhaustion now causes true mid-route stalling.\nDepartures are checked for fuel feasibility."),
        ("Port service realism", "Ports track actual vessel IDs.\nRefueling happens after real service completion."),
        ("Reward redesign", "Rewards moved from mostly aggregate penalties to event-driven, role-specific signals."),
        ("Final run cleanup", "Moved the finish line to continuous + ground_truth and fixed training resets to vary reproducibly."),
        ("Visibility", "Action logs, event logs, reward components, and confidence diagnostics were added."),
    ]
    positions = [(0.65, 1.45), (4.35, 1.45), (8.05, 1.45), (0.65, 4.15), (4.35, 4.15), (8.05, 4.15)]
    fills = [TEAL_SOFT, WHITE, ORANGE_SOFT, WHITE, TEAL_SOFT, WHITE]
    for (title, text), (x, y), fill in zip(items, positions, fills):
        _add_card(slide, x, y, 3.25, 2.15, fill=fill, title=title)
        _add_textbox(slide, x + 0.22, y + 0.6, 2.8, 1.15, text, font_size=13, color=INK)


def _diagnostics_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Why the current results are more trustworthy", "Diagnostics", "11")

    _add_card(slide, 0.65, 1.45, 4.35, 5.9, fill=WHITE, title="Transparency upgrades")
    box = _add_textbox(slide, 0.95, 1.95, 3.75, 4.9, "", font_size=1)
    _add_bullets(
        box,
        [
            "Grouped plots separate aggregate, vessel, port, and coordinator behavior.",
            "Per-step eval trace, action log, and event log are available.",
            "Reward-component decomposition shows which terms dominate.",
            "Policy-confidence metrics make near-uniform coordinator behavior easier to diagnose.",
        ],
        font_size=16,
    )

    image_path = _ensure_role_reward_plot()
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(5.45), Inches(1.85), width=Inches(6.85), height=Inches(3.1))
        _add_textbox(slide, 5.55, 5.0, 6.65, 0.32, "Role-specific reward traces over one evaluation episode. We show all three because vessel, port, and coordinator each optimize a different part of the task.", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
        chip_specs = [
            ("Trace file", "eval_trace.csv", 5.6, 5.55, 1.95),
            ("Action log", "eval_action_trace.csv", 7.75, 5.55, 2.15),
            ("Event log", "eval_event_log.csv", 10.1, 5.55, 1.95),
        ]
        for label, value, x, y, w in chip_specs:
            _add_card(slide, x, y, w, 0.92, fill=RGBColor(248, 253, 252))
            _add_textbox(slide, x + 0.12, y + 0.14, w - 0.24, 0.14, label.upper(), font_size=9, color=MUTED, bold=True)
            _add_textbox(slide, x + 0.12, y + 0.4, w - 0.24, 0.18, value, font_size=12, color=NAVY, bold=True)
    else:
        _add_card(slide, 5.35, 1.55, 7.3, 5.75, fill=RGBColor(245, 242, 236), title="Missing diagnostics image")


def _baseline_slide(prs: Presentation, v3: dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Current recommended baseline: transit-rebalanced v3", "Baseline", "12")

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.5), Inches(4.7), Inches(5.95))
    _style_shape(left, NAVY, None)
    _add_textbox(slide, 0.95, 1.85, 4.1, 0.32, "Preferred operating balance", font_size=20, color=PAPER, bold=True, font_name=FONT_HEAD)

    metrics = [
        ("Total reward", f"{v3['total_reward']:.2f}"),
        ("On-time rate", f"{v3['on_time_rate']:.3f}"),
        ("Completed arrivals", f"{v3['completed_arrivals']:.1f}"),
        ("Port service events", f"{v3['total_vessels_served']:.1f}"),
        ("Dock utilization", f"{v3['dock_utilization']:.2f}"),
        ("Ops cost", f"${v3['total_ops_cost_usd']/1_000_000:.3f}M"),
    ]
    positions = [(0.95, 2.45), (2.95, 2.45), (0.95, 3.85), (2.95, 3.85), (0.95, 5.25), (2.95, 5.25)]
    for (label, value), (x, y) in zip(metrics, positions):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.75), Inches(1.05))
        _style_shape(box, NAVY_2, None)
        _add_textbox(slide, x + 0.12, y + 0.12, 1.5, 0.14, label.upper(), font_size=9, color=SAND, bold=True)
        _add_textbox(slide, x + 0.12, y + 0.42, 1.5, 0.24, value, font_size=18, color=PAPER, bold=True, font_name=FONT_HEAD)

    image_path = V3_RUN / "training_curves.png"
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(5.75), Inches(1.55), width=Inches(6.7), height=Inches(5.75))
    else:
        _add_card(slide, 5.75, 1.55, 6.7, 5.75, fill=RGBColor(245, 242, 236), title="Missing training curves image")


def _benchmark_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "Final full-scale run plan", "Run plan", "13")
    panels = [
        ("Artifact run", 0.75, RGBColor(245, 250, 252), [
            ("Environment", "Continuous"),
            ("Forecast", "Ground truth"),
            ("Iterations", "100"),
            ("Seed", "42"),
            ("Purpose", "Figures, traces, report, and saved model"),
            ("Outputs", "report.md + trace CSVs + plots"),
        ]),
        ("Multi-seed run", 6.7, ORANGE_SOFT, [
            ("Environment", "Continuous"),
            ("Forecast", "Ground truth"),
            ("Iterations", "100"),
            ("Seeds", "42, 49, 56, 63, 70"),
            ("Purpose", "Final quantitative stability claim"),
            ("Outputs", "summary.csv + experiment_summary.json"),
        ]),
    ]
    for title, x, fill, stats in panels:
        _add_card(slide, x, 1.5, 5.85, 4.85, fill=fill, title=title)
        y = 2.05
        for idx, (label, value) in enumerate(stats):
            box_fill = WHITE if idx % 2 == 0 else RGBColor(250, 248, 244)
            _add_card(slide, x + 0.22, y, 5.4, 0.54, fill=box_fill)
            _add_textbox(slide, x + 0.36, y + 0.16, 3.75, 0.14, label.upper(), font_size=9, color=MUTED, bold=True)
            value_color = TEAL if title.startswith("Multi") and label in {"Seeds", "Outputs"} else NAVY
            _add_textbox(slide, x + 4.22, y + 0.14, 1.0, 0.16, value, font_size=14, color=value_color, bold=True, align=PP_ALIGN.RIGHT)
            y += 0.62

    _add_card(slide, 0.8, 6.55, 12.0, 0.95, fill=ORANGE_SOFT, title="Takeaway")
    _add_textbox(
        slide,
        1.05,
        6.94,
        11.45,
        0.28,
        "The project now finishes on one clean task definition: continuous scheduling with reliable forecasts, using one representative artifact run and one five-seed stability run.",
        font_size=15,
        color=INK,
    )


def _next_steps_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_title_band(slide, "What is simplified, and what comes next", "Next steps", "14")

    steps = [
        ("1. Run the artifact seed", "Generate the final figures, traces, report, and saved model from seed 42."),
        ("2. Run five seeds", "Use the same continuous + ground_truth setting for the final stability claim."),
        ("3. Write the result", "Use the artifact run for examples and the multi-seed run for the main numbers."),
        ("4. Keep scope honest", "Frame the result as control quality under reliable forecasts."),
        ("5. Future work", "Reintroduce imperfect forecasts and real-port geography afterward."),
    ]
    x = 0.75
    widths = [2.3, 2.3, 2.3, 2.3, 2.3]
    fills = [WHITE, TEAL_SOFT, WHITE, ORANGE_SOFT, WHITE]
    for (title, text), w, fill in zip(steps, widths, fills):
        _add_card(slide, x, 2.0, w, 3.0, fill=fill, title=title)
        _add_textbox(slide, x + 0.18, 2.6, w - 0.35, 1.7, text, font_size=14, color=INK)
        x += w + 0.18
    _add_textbox(slide, 0.95, 5.65, 11.4, 0.5, "This finish line is intentionally narrow: complete the project cleanly on the continuous environment, then expand realism afterward.", font_size=15, color=MUTED, align=PP_ALIGN.CENTER)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    v3 = _read_v3_metrics()
    _title_slide(prs)
    _context_slide(prs)
    _rq_slide(prs)
    _environment_slide(prs)
    _experiment_slide(prs)
    _architecture_slide(prs)
    _mechanics_slide(prs)
    _reward_summary_slide(prs)
    _reward_formula_slide(prs)
    _changes_slide(prs)
    _diagnostics_slide(prs)
    _baseline_slide(prs, v3)
    _benchmark_slide(prs)
    _next_steps_slide(prs)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_presentation())
